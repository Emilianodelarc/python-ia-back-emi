# import os, re, io, json
# from pathlib import Path
# from typing import List, Dict
# from pptx import Presentation

# OCR_ENABLED = os.getenv("OCR_ENABLED", "false").lower() == "true"
# try:
#     if OCR_ENABLED:
#         import pytesseract
#         from PIL import Image
#     else:
#         pytesseract = None
#         Image = None
# except Exception:
#     OCR_ENABLED = False
#     pytesseract = None
#     Image = None

# INPUT_DIR = "curso/ppt"
# OUT_JSONL = "curso/data/chunks.jsonl"
# CURSO = "PYTH_1200Funciones.pptx"


# def clean(t: str) -> str:
#     return re.sub(r"\s+", " ", (t or "").strip())


# def iter_shapes(shapes):
#     for sh in shapes:
#         yield sh
#         if hasattr(sh, "shapes"):
#             yield from iter_shapes(sh.shapes)


# def extract_table_text(shape) -> str:
#     try:
#         tbl = shape.table
#     except Exception:
#         return ""
#     out = []
#     for r in tbl.rows:
#         cells = [clean(c.text) for c in r.cells]
#         if any(cells):
#             out.append(" | ".join(cells))
#     return "\n".join(out)


# def extract_image_text(shape) -> str:
#     if not OCR_ENABLED or pytesseract is None or Image is None:
#         return ""
#     if not hasattr(shape, "image"):
#         return ""
#     try:
#         img_bytes = shape.image.blob
#         im = Image.open(io.BytesIO(img_bytes))
#         txt = pytesseract.image_to_string(im, lang="spa+eng")
#         return clean(txt)
#     except Exception:
#         return ""


# def extract_slide(slide) -> Dict:
#     parts = []
#     for sh in iter_shapes(slide.shapes):
#         if getattr(sh, "has_text_frame", False):
#             txt = clean(sh.text)
#             if txt:
#                 parts.append(txt)
#         # tablas
#         try:
#             if getattr(sh, "table", None) is not None:
#                 ttxt = extract_table_text(sh)
#                 if ttxt:
#                     parts.append(ttxt)
#         except Exception:
#             pass
#         # OCR en imágenes
#         ocr = extract_image_text(sh)
#         if ocr:
#             parts.append(f"[OCR] {ocr}")

#     # notas del orador
#     try:
#         if slide.has_notes_slide and slide.notes_slide and slide.notes_slide.notes_text_frame:
#             notes = clean(slide.notes_slide.notes_text_frame.text)
#             if notes:
#                 parts.append(f"[Notas] {notes}")
#     except Exception:
#         pass

#     full_text = "\n".join([p for p in parts if p])
#     return {"text": full_text}


# def chunk(text: str, max_chars=1000, overlap=120) -> List[str]:
#     text = clean(text)
#     if len(text) <= max_chars:
#         return [text]
#     chunks = []
#     i = 0
#     while i < len(text):
#         chunks.append(text[i:i + max_chars])
#         i = i + max_chars - overlap
#     return chunks


# def run_folder(input_dir=INPUT_DIR, out_jsonl=OUT_JSONL, curso=CURSO):
#     Path(out_jsonl).parent.mkdir(parents=True, exist_ok=True)
#     with open(out_jsonl, "w", encoding="utf-8") as f:
#         for pptx in sorted(Path(input_dir).glob("*.pptx")):
#             clase = pptx.stem  # p.ej., clase01
#             prs = Presentation(str(pptx))
#             for sidx, slide in enumerate(prs.slides, start=1):
#                 data = extract_slide(slide)
#                 if not data["text"]:
#                     continue
#                 for ci, ch in enumerate(chunk(data["text"]), start=1):
#                     row = {
#                         "curso": curso,
#                         "clase": clase,
#                         "slide": sidx,
#                         "chunk_id": f"{clase}-s{sidx}-c{ci}",
#                         "text": ch,
#                     }
#                     f.write(json.dumps(row, ensure_ascii=False) + "\n")
#     print(f"Listo: {out_jsonl}")


# if __name__ == "__main__":
#     run_folder()

import os, re, io, json
from pathlib import Path
from typing import List, Dict
from pptx import Presentation

DEBUG = True  # 👈 ponlo en False cuando todo ande

OCR_ENABLED = os.getenv("OCR_ENABLED", "false").lower() == "true"
try:
    if OCR_ENABLED:
        import pytesseract
        from PIL import Image
    else:
        pytesseract = None
        Image = None
except Exception:
    OCR_ENABLED = False
    pytesseract = None
    Image = None

INPUT_DIR = "curso/ppt"
OUT_JSONL = "curso/data/chunks.jsonl"

def clean(t: str) -> str:
    return re.sub(r"\s+", " ", (t or "").strip())

def iter_shapes(shapes):
    for sh in shapes:
        yield sh
        if hasattr(sh, "shapes"):
            yield from iter_shapes(sh.shapes)

def extract_table_text(shape) -> str:
    try:
        tbl = shape.table
    except Exception:
        return ""
    out = []
    for r in tbl.rows:
        cells = [clean(c.text) for c in r.cells]
        if any(cells):
            out.append(" | ".join(cells))
    return "\n".join(out)

def extract_image_text(shape) -> str:
    if not OCR_ENABLED or pytesseract is None or Image is None:
        return ""
    if not hasattr(shape, "image"):
        return ""
    try:
        img_bytes = shape.image.blob
        im = Image.open(io.BytesIO(img_bytes))
        txt = pytesseract.image_to_string(im, lang="spa+eng")
        return clean(txt)
    except Exception:
        return ""

def extract_slide(slide) -> Dict:
    parts = []
    for sh in iter_shapes(slide.shapes):
        if getattr(sh, "has_text_frame", False):
            txt = clean(sh.text)
            if txt:
                parts.append(txt)
        try:
            if getattr(sh, "table", None) is not None:
                ttxt = extract_table_text(sh)
                if ttxt:
                    parts.append(ttxt)
        except Exception:
            pass
        ocr = extract_image_text(sh)
        if ocr:
            parts.append(f"[OCR] {ocr}")
    try:
        if slide.has_notes_slide and slide.notes_slide and slide.notes_slide.notes_text_frame:
            notes = clean(slide.notes_slide.notes_text_frame.text)
            if notes:
                parts.append(f"[Notas] {notes}")
    except Exception:
        pass
    full_text = "\n".join([p for p in parts if p])
    return {"text": full_text}

def chunk(text: str, max_chars=1000, overlap=120) -> List[str]:
    text = clean(text)
    if len(text) <= max_chars:
        return [text]
    chunks = []
    i = 0
    while i < len(text):
        chunks.append(text[i:i + max_chars])
        i = i + max_chars - overlap
    return chunks

def run_folder(input_dir=INPUT_DIR, out_jsonl=OUT_JSONL):
    Path(out_jsonl).parent.mkdir(parents=True, exist_ok=True)
    pptx_files = sorted(Path(input_dir).glob("*.pptx"))
    if DEBUG:
        print(f"[DEBUG] cwd = {Path().resolve()}")
        print(f"[DEBUG] INPUT_DIR exists = {Path(input_dir).exists()}")
        print(f"[DEBUG] PPTX encontrados = {len(pptx_files)}")
        for p in pptx_files[:10]:
            print(" -", p.name)

    total_chunks = 0
    total_slides = 0
    total_with_text = 0

    with open(out_jsonl, "w", encoding="utf-8") as f:
        for pptx in pptx_files:
            curso = pptx.name
            clase = pptx.stem
            try:
                prs = Presentation(str(pptx))
            except Exception as e:
                print(f"[WARN] No pude abrir {pptx.name}: {e}")
                continue

            if DEBUG:
                print(f"\n[DEBUG] Procesando: {pptx.name} | slides={len(prs.slides)}")

            for sidx, slide in enumerate(prs.slides, start=1):
                total_slides += 1
                data = extract_slide(slide)
                if not data["text"]:
                    if DEBUG:
                        print(f"  - Slide {sidx}: (sin texto)")
                    continue

                total_with_text += 1
                chs = chunk(data["text"])
                if DEBUG:
                    print(f"  - Slide {sidx}: {len(data['text'])} chars → {len(chs)} chunks")

                for ci, ch in enumerate(chs, start=1):
                    row = {
                        "curso": curso,
                        "clase": clase,
                        "slide": sidx,
                        "chunk_id": f"{clase}-s{sidx}-c{ci}",
                        "text": ch,
                    }
                    f.write(json.dumps(row, ensure_ascii=False) + "\n")
                    total_chunks += 1

    print(f"\nListo: {out_jsonl}")
    print(f"Slides totales: {total_slides}")
    print(f"Slides con texto: {total_with_text}")
    print(f"Chunks escritos: {total_chunks}")

if __name__ == "__main__":
    run_folder()
